import logging
import os
import json
from openai import OpenAI
from django.conf import settings
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import os
import logging
import json
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from django.conf import settings
from openai import OpenAI
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
from io import BytesIO

logger = logging.getLogger(__name__)
client = OpenAI(api_key=getattr(settings, "OPENAI_API_KEY", None))

import os
import logging
import traceback
from openai import OpenAI
from django.conf import settings

logger = logging.getLogger(__name__)
client = OpenAI(api_key=settings.OPENAI_API_KEY)


from openai import OpenAI, BadRequestError
from django.conf import settings
import logging

logger = logging.getLogger(__name__)
client = OpenAI(api_key=settings.OPENAI_API_KEY)

import mimetypes

from pdf2image import convert_from_bytes
from mimetypes import guess_type

def summarize_file_with_vision(file_obj):
    path = file_obj.file.path
    mime, _ = guess_type(path)
    # Decide route based on MIME
    if mime in ("image/png", "image/jpeg", "image/webp", "image/gif"):
        return _summarize_image(path)
    elif mime == "application/pdf":
        return _summarize_pdf(path)
    else:
        # Extract text fallback
        text = extract_text_from_file(file_obj)
        if not text.strip():
            return None, "Unsupported file type and text extraction failed."
        return summarize_text(text), None


def _summarize_image(path):
    with open(path, "rb") as f:
        upload = client.files.create(file=f, purpose="assistants")
    return _vision_chat(upload.id)

def _summarize_pdf(path):
    images = convert_from_bytes(open(path, "rb").read(), fmt="PNG")
    summaries, warnings = [], []
    for img in images:
        buf = BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        upload = client.files.create(file=buf, purpose="assistants")
        summary, warn = _vision_chat(upload.id)
        if summary: summaries.append(summary)
        if warn: warnings.append(warn)
    return "\n\n".join(summaries), "; ".join(warnings) or None

def _vision_chat(file_id):
    resp = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": [
                {"type": "text", "text": "Summarize this image."},
                {"type": "image_url", "image_url": {"url": f"file-id:{file_id}"}}
            ]}
        ],
        temperature=0.4, max_tokens=600
    )
    text = resp.choices[0].message.content.strip()
    warning = None
    if len(text.split()) >= 450:
        warning = "Might be partial due to token limit."
    return text, warning

def summarize_text(text):
    resp = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": f"Summarize this text:\n\n{ text[:3000] }"}
        ],
        temperature=0.4, max_tokens=600
    )
    return resp.choices[0].message.content.strip()



def transcribe_audio(file_path: str) -> str:
    """
    Transcribes audio file using Whisper API.
    """
    try:
        with open(file_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                file=audio_file,
                model="whisper-1",
                response_format="text"
            )
        return transcript.strip()
    except Exception as e:
        logger.warning(f"OpenAI transcription failed: {e}")
        return "Transcription failed."


def extract_text_from_file(file_obj) -> str:
    """
    Extracts text from PDF, DOCX, or PPTX file.
    """
    try:
        file_path = file_obj.path
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)

        elif ext == ".docx":
            doc = Document(file_path)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs)

        elif ext == ".pptx":
            ppt = Presentation(file_path)
            text_runs = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        else:
            raise ValueError("Unsupported file type for extraction.")

    except Exception as e:
        logger.warning(f"Failed to extract text: {e}")
        return ""
# files/utils.py

def generate_questions_from_text(
    text: str,
    mode: str = "both",
    visuals: bool = False,
    difficulty: str = "medium",
    num_questions: int = 5
) -> tuple[list[dict], str | None]:
    """
    Generates questions from text using GPT-4o.

    Returns:
        - A list of parsed question dicts.
        - A warning message (or None if all went well).
    """
    try:
        instructions = f"Generate {num_questions} questions from the given study material.\n"

        if mode == "mcq":
            instructions += "Only generate multiple-choice questions (with 4 options).\n"
        elif mode == "theory":
            instructions += "Only generate open-ended theory questions.\n"
        else:
            instructions += "Include a mix of MCQs and theory questions.\n"

        if visuals:
            instructions += "at least 1 must include visual hints (like describing an image or asking for a graph).\n"

        instructions += f"Difficulty level should be '{difficulty}'. Return only a valid JSON array in the format below:\n"
        instructions += """
        [
            {
                "question": "...",
                "answer": "...",
                "type": "mcq" or "theory",
                "options": [optional if MCQ],
                "visual_aid": [optional if visual],
                "difficulty": "easy/medium/hard"
            }
        ]
        """

        prompt = instructions + "\n\nText:\n" + text[:6000]

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=1000
        )

        raw = response.choices[0].message.content.strip()
        logger.info("\n==== RAW OPENAI OUTPUT ====\n" + raw + "\n===========================\n")

        # Always parse safely and return both questions and warning
        questions, warning = parse_openai_question_output(
            raw, expected=num_questions, difficulty=difficulty
        )
        return questions, warning

    except Exception as e:
        logger.warning(f"OpenAI question generation failed: {e}")
        return [], "Something went wrong while generating the questions. Please try again."

logger = logging.getLogger(__name__)
client = OpenAI(api_key=getattr(settings, "OPENAI_API_KEY", None))


def extract_text_and_images_from_file(file_obj):
    path = file_obj.path
    ext = os.path.splitext(path)[1].lower()
    images = {}
    text = ""

    if ext == ".pdf":
        reader = PdfReader(path)
        for i, page in enumerate(reader.pages):
            text += page.extract_text() or ""

        for page_index, page in enumerate(reader.pages):
            try:
                xObject = page["/Resources"].get("/XObject", {})
                xObject = xObject.get_object() if hasattr(xObject, "get_object") else xObject
                for obj in xObject:
                    if xObject[obj].get("/Subtype") == "/Image":
                        img_data = xObject[obj]._data
                        image = Image.open(BytesIO(img_data)).convert("RGB")
                        filename = f"visual_qs/pdf_img_{page_index}_{obj[1:]}.jpg"
                        image_io = BytesIO()
                        image.save(image_io, format="JPEG")
                        saved_path = default_storage.save(filename, ContentFile(image_io.getvalue()))
                        images[f"image_{len(images)+1}"] = default_storage.url(saved_path)
            except Exception as e:
                logger.warning(f"PDF image extract failed: {e}")

    elif ext == ".docx":
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs)

        rels = doc.part._rels
        for rel in rels:
            rel = rels[rel]
            if "image" in rel.target_ref:
                try:
                    img_blob = rel.target_part.blob
                    image = Image.open(BytesIO(img_blob)).convert("RGB")
                    filename = f"visual_qs/docx_img_{len(images)}.jpg"
                    image_io = BytesIO()
                    image.save(image_io, format="JPEG")
                    saved_path = default_storage.save(filename, ContentFile(image_io.getvalue()))
                    images[f"image_{len(images)+1}"] = default_storage.url(saved_path)
                except Exception as e:
                    logger.warning(f"DOCX image extract failed: {e}")

    elif ext == ".pptx":
        ppt = Presentation(path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                if shape.shape_type == 13:
                    try:
                        img_blob = shape.image.blob
                        image = Image.open(BytesIO(img_blob)).convert("RGB")
                        filename = f"visual_qs/pptx_img_{len(images)}.jpg"
                        image_io = BytesIO()
                        image.save(image_io, format="JPEG")
                        saved_path = default_storage.save(filename, ContentFile(image_io.getvalue()))
                        images[f"image_{len(images)+1}"] = default_storage.url(saved_path)
                    except Exception as e:
                        logger.warning(f"PPTX image extract failed: {e}")

    return text.strip(), images


def parse_openai_question_output(raw: str, expected: int = 5, difficulty: str = "medium"):
    warning = None
    if raw.startswith("```"):
        raw = raw.strip("`").strip()
        if raw.lower().startswith("json"):
            raw = raw[4:].strip()

    parsed_questions = []
    try:
        parsed = json.loads(raw)
        if isinstance(parsed, list):
            parsed_questions = [
                q for q in parsed
                if isinstance(q, dict) and q.get("question") and q.get("answer")
            ]
    except json.JSONDecodeError as e:
        logger.warning(f"JSON parsing failed: {e}")
        logger.debug("Failed JSON:\n" + raw[:300])

    if not parsed_questions:
        for block in raw.split("\n\n"):
            if "Q:" in block and "A:" in block:
                try:
                    q = block.split("Q:")[1].split("A:")[0].strip()
                    a = block.split("A:")[1].strip()
                    parsed_questions.append({
                        "question": q,
                        "answer": a,
                        "type": "theory",
                        "difficulty": difficulty
                    })
                except Exception as fallback_error:
                    logger.warning(f"Fallback parse error: {fallback_error}")
                    continue

    if len(parsed_questions) < expected:
        warning = f"Only {len(parsed_questions)} out of {expected} questions were generated."

    return parsed_questions, warning


def generate_questions_from_text_and_images(
    text: str,
    image_map: dict[str, str],
    mode: str = "both",
    visuals: bool = False,
    difficulty: str = "medium",
    num_questions: int = 5
):
    try:
        image_refs = "\n".join([f"{k}: {v}" for k, v in image_map.items()])
        instructions = f"Generate {num_questions} exam questions. "

        if mode == "mcq":
            instructions += "Use only multiple choice format. "
        elif mode == "theory":
            instructions += "Use only theory (open-ended) questions. "
        else:
            instructions += "Use a mix of MCQs and theory. "

        if visuals and image_map:
            instructions += (
                "at least 1 question must refer to the provided images (e.g. 'Refer to image_2'). "
            )

        instructions += f"Difficulty: {difficulty}. Format output as a JSON list with keys: question, answer, type, difficulty, and optionally options and visual_aid."

        prompt = f"""{instructions}

Images:
{image_refs}

Text:
{text[:6000]}
"""

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            max_tokens=1500
        )

        raw = response.choices[0].message.content.strip()
        return parse_openai_question_output(raw, expected=num_questions, difficulty=difficulty)

    except Exception as e:
        logger.warning(f"OpenAI question gen with visuals failed: {e}")
        return [], "Something went wrong while generating questions."

def score_theory_answer(question_text: str, correct_answer: str, user_answer: str) -> dict:
    """
    Uses GPT-4o to semantically score a user's theory answer against the correct answer.
    Returns a score (0.0 to 1.0) and a short justification.
    """
    try:
        prompt = f"""
Evaluate the following student answer.

- Compare it against the expected answer.
- Score it from 0 (completely wrong) to 1 (fully correct).
- Also explain why you gave that score.

Return ONLY JSON like:
{{ "score": float (0 to 1), "comment": "brief justification" }}

Question: {question_text}
Expected Answer: {correct_answer}
Student Answer: {user_answer}
"""

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an educational assistant scoring theory questions."},
                {"role": "user", "content": prompt}
            ],
            temperature=0,
            max_tokens=300
        )

        content = response.choices[0].message.content.strip()

        # Try parsing JSON
        if content.startswith("```"):
            content = content.strip("`").strip()
            if content.lower().startswith("json"):
                content = content[4:].strip()

        result = json.loads(content)

        return {
            "score": float(result.get("score", 0.0)),
            "comment": result.get("comment", "No comment.")
        }

    except Exception as e:
        logger.warning(f"Theory scoring failed: {e}")
        return {
            "score": 0.0,
            "comment": "Evaluation failed. Please try again."
        }
